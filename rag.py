import os
import requests
import faiss
import pickle
from tqdm import tqdm
from typing import List, Dict, Any
from pathlib import Path
import numpy as np

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    import pptx
except ImportError:
    pptx = None

INDEX_PATH = "rag.index"
META_PATH = "rag.meta.pkl"
MISTRAL_API_URL = "https://api.mistral.ai/v1/embeddings"
MISTRAL_MODEL = "mistral-embed"
MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY")

RACINE = "OneDrive_1_02-07-2025"
AO_TYPES = {"AO_GO": "GO", "AO_NOGO": "NOGO"}
EXTENSIONS = ("pdf", "docx", "txt", "csv", "xls", "xlsx", "pptx")

# Extraction de texte selon le type de fichier
def extract_text_from_file(filepath: str) -> str:
    ext = filepath.lower().split(".")[-1]
    if ext == "pdf" and PdfReader:
        try:
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            print(f"Erreur lecture PDF {filepath}: {e}")
            return ""
    elif ext == "docx" and docx:
        try:
            doc = docx.Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            print(f"Erreur lecture DOCX {filepath}: {e}")
            return ""
    elif ext == "txt":
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"Erreur lecture TXT {filepath}: {e}")
            return ""
    elif ext == "csv" and pd:
        try:
            df = pd.read_csv(filepath, dtype=str, encoding_errors='ignore')
            return df.astype(str).to_string(index=False)
        except Exception as e:
            print(f"Erreur lecture CSV {filepath}: {e}")
            return ""
    elif ext in ("xls", "xlsx") and pd:
        try:
            if ext == "xls":
                try:
                    df = pd.read_excel(filepath, dtype=str, engine=None)
                except ImportError:
                    print("xlrd n'est pas installé : pip install xlrd pour lire les .xls")
                    return ""
            else:
                df = pd.read_excel(filepath, dtype=str, engine='openpyxl')
            return df.astype(str).to_string(index=False)
        except Exception as e:
            print(f"Erreur lecture Excel {filepath}: {e}")
            return ""
    elif ext == "pptx" and pptx:
        try:
            prs = pptx.Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            print(f"Erreur lecture PPTX {filepath}: {e}")
            return ""
    else:
        return ""

def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 200) -> List[str]:
    # Découpe le texte en chunks de taille chunk_size avec un overlap
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def get_all_documents() -> List[Dict[str, Any]]:
    docs = []
    racine = Path(RACINE)
    for ao_folder, decision in AO_TYPES.items():
        ao_dir = racine / ao_folder
        if not ao_dir.exists():
            continue
        for ao in ao_dir.iterdir():
            if not ao.is_dir():
                continue
            ao_id = ao.name
            for ext in EXTENSIONS:
                for filepath in ao.rglob(f"*.{ext}"):
                    rel_path = filepath.relative_to(racine)
                    text = extract_text_from_file(str(filepath))
                    if text.strip():
                        for chunk in chunk_text(text):
                            docs.append({
                                "content": chunk,
                                "metadata": {
                                    "filepath": str(rel_path),
                                    "ao_id": ao_id,
                                    "decision": decision
                                }
                            })
    return docs

def get_mistral_embedding(texts: List[str]) -> List[List[float]]:
    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": MISTRAL_MODEL,
        "input": texts
    }
    try:
        response = requests.post(MISTRAL_API_URL, headers=headers, json=data, timeout=30)
        response.raise_for_status()
        return [item["embedding"] for item in response.json()["data"]]
    except Exception as e:
        print(f"Erreur API Mistral: {e}")
        return [[0.0]*1024 for _ in texts]  # vecteur nul si erreur

def build_index(documents: List[Dict[str, Any]]):
    print(f"Vectorisation de {len(documents)} chunks...")
    embeddings = []
    metadatas = []
    for i in tqdm(range(0, len(documents), 32)):
        batch = documents[i:i+32]
        texts = [doc["content"] for doc in batch]
        batch_embeds = get_mistral_embedding(texts)
        embeddings.extend(batch_embeds)
        metadatas.extend([
            {
                **doc["metadata"],
                "content": doc["content"]
            } for doc in batch
        ])
    dim = len(embeddings[0])
    index = faiss.IndexFlatL2(dim)
    index.add(np.array(embeddings).astype('float32'))
    save_index(index, metadatas)
    print(f"Index sauvegardé dans {INDEX_PATH} et {META_PATH}")

def save_index(index, metadatas):
    faiss.write_index(index, INDEX_PATH)
    with open(META_PATH, "wb") as f:
        pickle.dump(metadatas, f)

def load_index():
    index = faiss.read_index(INDEX_PATH)
    with open(META_PATH, "rb") as f:
        metadatas = pickle.load(f)
    return index, metadatas

def query_rag(question: str, top_k: int = 5):
    index, metadatas = load_index()
    embed = get_mistral_embedding([question])[0]
    D, I = index.search(np.array([embed]).astype('float32'), top_k)
    results = []
    for score, idx in zip(D[0], I[0]):
        meta = metadatas[idx]
        results.append({
            "score": float(1/(1+score)),
            "content": meta.get("content", "[non disponible]"),
            "filepath": meta["filepath"],
            "ao_id": meta["ao_id"],
            "decision": meta["decision"]
        })
    return results

if __name__ == "__main__":
    # MODE TEST : ne vectoriser que 2 fichiers pour valider la pipeline
    docs = []
    racine = Path(RACINE)
    count = 0
    for ao_folder, decision in AO_TYPES.items():
        ao_dir = racine / ao_folder
        if not ao_dir.exists():
            continue
        for ao in ao_dir.iterdir():
            if not ao.is_dir():
                continue
            ao_id = ao.name
            for ext in EXTENSIONS:
                for filepath in ao.rglob(f"*.{ext}"):
                    rel_path = filepath.relative_to(racine)
                    text = extract_text_from_file(str(filepath))
                    if text.strip():
                        for chunk in chunk_text(text):
                            docs.append({
                                "content": chunk,
                                "metadata": {
                                    "filepath": str(rel_path),
                                    "ao_id": ao_id,
                                    "decision": decision
                                }
                            })
                        count += 1
                        if count >= 2:
                            break
                if count >= 2:
                    break
            if count >= 2:
                break
        if count >= 2:
            break
    print(f"Test : {len(docs)} chunks extraits à partir de 2 fichiers.")
    build_index(docs)
    print("Test terminé. Index construit sur 2 fichiers.") 