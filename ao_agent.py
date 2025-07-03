import uuid
from typing import Any
from langchain_core.messages import HumanMessage
from langchain_core.tools import tool
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import create_react_agent
import os
from mistralai import Mistral
from langchain_core.language_models.chat_models import SimpleChatModel
from langchain_core.messages.ai import AIMessage
from pydantic import Field
import logging

# Configuration logging globale
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.FileHandler("agent_full.log", mode="w", encoding="utf-8"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("agent")
logger.setLevel(logging.DEBUG)

# Outil : query_rag
@tool
def query_rag_tool(question: str) -> Any:
    """Récupère des extraits pertinents d'AO passés (Go/NoGo) via RAG."""
    from rag import query_rag
    return query_rag(question)

# 1. ShortMemory
class ShortMemory(MemorySaver):
    def put(self, config, state, *args, **kwargs):
        # Ne garde que les 2 derniers messages dans la mémoire
        if "messages" in state and isinstance(state["messages"], list):
            state = dict(state)  # shallow copy
            state["messages"] = state["messages"][-2:]
        return super().put(config, state, *args, **kwargs)

# Outil : read_documents_from_folder (ne lit que le début de chaque document)
@tool
def read_documents_from_folder_tool(folder_path: str, max_chars: int = 1500) -> Any:
    """Concatène le début (≤ 1 500 car.) de chaque fichier du dossier AO."""
    import glob
    from pathlib import Path
    import docx, pandas as pd
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None
    try:
        import openpyxl
    except ImportError:
        openpyxl = None
    try:
        import pptx
    except ImportError:
        pptx = None
    exts = ("pdf", "doc", "docx", "xls", "xlsx", "txt", "csv", "pptx")
    texts = []
    long_files = []
    for ext in exts:
        for file in Path(folder_path).rglob(f"*.{ext}"):
            try:
                content = ""
                if ext == "pdf" and PdfReader:
                    with open(file, "rb") as f:
                        reader = PdfReader(f)
                        content = "\n".join(page.extract_text() or "" for page in reader.pages)
                elif ext == "docx":
                    doc = docx.Document(file)
                    content = "\n".join([p.text for p in doc.paragraphs])
                elif ext == "txt":
                    with open(file, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                elif ext == "csv":
                    df = pd.read_csv(file, dtype=str, encoding_errors='ignore')
                    content = df.astype(str).to_string(index=False)
                elif ext in ("xls", "xlsx"):
                    df = pd.read_excel(file, dtype=str, engine='openpyxl' if ext=="xlsx" else None)
                    content = df.astype(str).to_string(index=False)
                elif ext == "pptx" and pptx:
                    prs = pptx.Presentation(file)
                    content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
                if content:
                    texts.append(f"--- {file} ---\n" + content[:max_chars])
                    if len(content) > max_chars:
                        long_files.append({
                            "filepath": str(file),
                            "length": len(content),
                            "has_more": True
                        })
            except Exception as e:
                continue
    return {
        "text": "\n\n".join(texts),
        "long_files": long_files
    }

# Outil : read_more_from_file
@tool
def read_more_from_file(filepath: str, offset: int = 4000, max_chars: int = 4000) -> str:
    """Lit un segment supplémentaire d'un fichier long (pagination)."""
    from pathlib import Path
    import docx, pandas as pd
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None
    try:
        import openpyxl
    except ImportError:
        openpyxl = None
    try:
        import pptx
    except ImportError:
        pptx = None
    ext = Path(filepath).suffix.lower().replace('.', '')
    content = ""
    try:
        if ext == "pdf" and PdfReader:
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == "docx":
            doc = docx.Document(filepath)
            content = "\n".join([p.text for p in doc.paragraphs])
        elif ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        elif ext == "csv":
            df = pd.read_csv(filepath, dtype=str, encoding_errors='ignore')
            content = df.astype(str).to_string(index=False)
        elif ext in ("xls", "xlsx"):
            df = pd.read_excel(filepath, dtype=str, engine='openpyxl' if ext=="xlsx" else None)
            content = df.astype(str).to_string(index=False)
        elif ext == "pptx" and pptx:
            prs = pptx.Presentation(filepath)
            content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e:
        return f"Erreur lors de la lecture de {filepath} : {e}"
    if not content:
        return f"Aucun contenu lisible dans {filepath}"
    return content[offset:offset+max_chars]

# Nouveau SYSTEM_PROMPT pour orchestration autonome
SYSTEM_PROMPT = """
Tu es consultant data. À partir de documents AO et du RAG, produis :
1) résumé, 2) points à creuser, 3) Go/NoGo argumenté, 4) plan de réponse, 5) synthèse manager.
Structure toujours ta réponse en 5 blocs Markdown nommés comme ci-dessous.
"""

def make_agent():
    memory = ShortMemory()
    mistral_api_key = os.getenv("MISTRAL_API_KEY", "")
    mistral_model = "mistral-large-latest"
    class MistralChatModel(SimpleChatModel):
        api_key: str = Field(default="", exclude=True)
        model: str = Field(default="mistral-small-latest", exclude=True)
        lc_tools: list = Field(default_factory=list, exclude=True)

        def __init__(self, api_key, model, lc_tools=None, **kwargs):
            super().__init__(**kwargs)
            object.__setattr__(self, "api_key", api_key)
            object.__setattr__(self, "model", model)
            object.__setattr__(self, "lc_tools", lc_tools or [])

        @property
        def _llm_type(self):
            return "mistral-chat"

        def bind_tools(self, tools, tool_choice=None):
            return self.__class__(
                api_key=self.api_key,
                model=self.model,
                lc_tools=tools,
                **self.dict(exclude={"api_key", "model", "lc_tools"})
            )

        def _call(self, messages, stop=None, run_manager=None, **kwargs):
            # Conversion des messages LangChain en format Mistral
            mistral_messages = []
            for m in messages:
                if hasattr(m, "type"):
                    if m.type == "human":
                        mistral_messages.append({"role": "user", "content": m.content})
                    elif m.type == "ai":
                        # On ignore les messages assistant pour Mistral (sinon erreur 400)
                        continue
                    elif m.type == "tool":
                        mistral_messages.append({"role": "tool", "content": m.content})
                    elif m.type == "system":
                        mistral_messages.append({"role": "system", "content": m.content})
                else:
                    # fallback : on traite comme user
                    mistral_messages.append({"role": "user", "content": getattr(m, "content", str(m))})
            # S'assurer que le dernier message est bien user ou tool
            while mistral_messages and mistral_messages[-1]["role"] not in ("user", "tool"):
                mistral_messages.pop()
            with Mistral(api_key=self.api_key) as mistral:
                res = mistral.chat.complete(
                    model=self.model,
                    messages=mistral_messages,
                )
            return res.choices[0].message.content
    model = MistralChatModel(mistral_api_key, mistral_model)
    tools = [read_documents_from_folder_tool, read_more_from_file, query_rag_tool]
    agent = create_react_agent(
        model,
        tools=tools,
        prompt=SYSTEM_PROMPT,
        checkpointer=memory,
    )
    return agent, memory

def print_agent_action(event):
    """Affiche l'action de l'agent (lecture, appel rag, etc.), le content des réponses intermédiaires, et les fichiers lus via RAG. Loggue tout dans agent_full.log."""
    if "tool" in event:
        if event["tool"] == "read_documents_from_folder_tool":
            msg = f"[Agent] Lecture des documents du dossier AO..."
            print(msg)
            logger.info(msg)
        elif event["tool"] == "query_rag_tool":
            msg = f"[Agent] Appel à query_rag_tool : {event.get('tool_input','')[:60]}..."
            print(msg)
            logger.info(msg)
            result = event.get("output", None)
            if isinstance(result, list):
                filepaths = [r.get("filepath") for r in result if isinstance(r, dict) and "filepath" in r]
                if filepaths:
                    print("[Agent] Fichiers consultés via RAG :")
                    logger.info("[Agent] Fichiers consultés via RAG :")
                    for fp in filepaths:
                        print(f"   - {fp}")
                        logger.info(f"   - {fp}")
        elif event["tool"] == "read_more_from_file":
            msg = f"[Agent] Lecture d'un extrait supplémentaire : {event.get('tool_input','')[:60]}..."
            print(msg)
            logger.info(msg)
        else:
            msg = f"[Agent] Utilisation de l'outil : {event['tool']}"
            print(msg)
            logger.info(msg)
    elif "messages" in event:
        msg = event["messages"][-1]
        if hasattr(msg, 'type') and msg.type == 'human':
            print("[Utilisateur] Nouvelle consigne envoyée à l'agent.")
            logger.info("[Utilisateur] Nouvelle consigne envoyée à l'agent.")
        elif hasattr(msg, 'type') and msg.type == 'ai':
            print("[Agent] Génération d'une réponse intermédiaire :\n")
            logger.info("[Agent] Génération d'une réponse intermédiaire :\n")
            print(msg.content)
            logger.debug(msg.content)

def run_agent_on_folder(folder_path: str):
    agent, memory = make_agent()
    thread_id = uuid.uuid4()
    config = {"configurable": {"thread_id": thread_id}}
    # 1. Lecture et résumé initial (préparation du contexte)
    docs_info = read_documents_from_folder_tool.invoke(folder_path)
    if isinstance(docs_info, dict):
        text = docs_info.get("text", "")
        long_files = docs_info.get("long_files", [])
    else:
        text = docs_info
        long_files = []
    # Limite stricte : max 1200 caractères au total, 300 par doc, max 3 docs
    MAX_DOCS = 3
    DOC_SLICE = 300
    TOTAL_MAX = 1200
    docs_split = text.split("--- ")
    selected_docs = []
    char_count = 0
    for doc in docs_split:
        if len(selected_docs) >= MAX_DOCS:
            break
        doc = doc.strip()
        if not doc:
            continue
        doc = doc[:DOC_SLICE]
        if char_count + len(doc) > TOTAL_MAX:
            break
        selected_docs.append(doc)
        char_count += len(doc)
    final_text = "\n--- ".join(selected_docs)
    if len(text) > len(final_text):
        final_text += "\n\n[Texte tronqué pour respecter la limite stricte de contexte]"
    input_message = HumanMessage(content=f"Voici le dossier AO à analyser :\n\n{final_text}")

    # 2. Exploration RAG sur des thèmes clés (préparation du contexte RAG)
    themes = [
        "procédure avec négociation",
        "budget estimé",
        "plateforme Datalake",
        "durée du contrat",
        "clauses RGPD"
    ]
    rag_insights = []
    for theme in themes:
        result = query_rag_tool.invoke(theme)
        if isinstance(result, list):
            result = result[:2]
        rag_insights.append(f"**{theme}**:\n{str(result)[:1200]}\n")
    rag_summary = "\n\n".join(rag_insights)
    rag_summary = rag_summary[:6000]
    rag_message = HumanMessage(content=f"Voici les résultats RAG à intégrer à ton analyse :\n\n{rag_summary}")

    # 3. Lancement de l'agent : une seule génération finale, avec tout le contexte
    print("\n--- Analyse complète (Go/NoGo, plan, synthèse manager) ---\n")
    # On envoie les deux messages dans la même conversation
    messages = [input_message, rag_message]
    for event in agent.stream({"messages": messages}, config, stream_mode="values"):
        print_agent_action(event)

if __name__ == "__main__":
    folder_path = input("Chemin du dossier AO à analyser : ")
    run_agent_on_folder(folder_path) 